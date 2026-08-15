#!/usr/bin/env python3
"""Create a folder-preserving, pseudonymized copy of an attachment archive.

This utility is deliberately conservative: it never edits input files and it
reports every skipped file. The reverse mapping is always written as plain,
readable JSON.
"""
from __future__ import annotations

import argparse
import csv
import hashlib
import hmac
import json
import os
import re
import shutil
from collections import Counter
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable

EMAIL = re.compile(r"\b[A-Za-z0-9.!#$%&'*+/=?^_`{|}~-]+@[A-Za-z0-9-]+(?:\.[A-Za-z0-9-]+)+\b")
PHONE = re.compile(r"(?<!\w)(?:\+?\d{1,3}[ .-]?)?(?:\(?\d{2,4}\)?[ .-]?)?\d{3}[ .-]\d{3,4}(?!\w)")
SIN = re.compile(r"(?<!\d)\d{3}[ -]?\d{3}[ -]?\d{3}(?!\d)")
CARD = re.compile(r"(?<!\d)(?:\d[ -]*?){13,19}(?!\d)")
IP = re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b")
COMPANY_CUE = re.compile(
    r"\b(?:inc\.?|incorporated|ltd\.?|limited|llc|llp|plc|corp\.?|corporation|"
    r"company|co\.?|bank|holdings?|group|partners?|trust|funds?)\b",
    re.IGNORECASE,
)
SUPPORTED = {".docx", ".xlsx", ".xlsm", ".pptx", ".pdf", ".csv", ".tsv", ".txt", ".md", ".json", ".xml", ".html", ".htm", ".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp"}
IMAGE_TYPES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".webp"}


def normalise(value: str) -> str:
    return re.sub(r"[^a-z0-9]", "", value.casefold())


def _word_tokens(value: str) -> list[str]:
    return re.findall(r"[a-z0-9]+", value.casefold())


def _shares_full_word_span(a: list[str], b: list[str]) -> bool:
    """True if the shorter word list appears as a contiguous run inside the longer one."""
    if not a or not b:
        return False
    longer, shorter = (a, b) if len(a) >= len(b) else (b, a)
    for i in range(len(longer) - len(shorter) + 1):
        if longer[i:i + len(shorter)] == shorter:
            return True
    return False


CODE_LETTER = {
    "company": "C", "person": "P", "phone": "T",
    "id": "I", "card": "K", "ip": "N", "label": "G",
}


@dataclass
class Pseudonymizer:
    secret: bytes
    configured_entities: dict[str, str] = field(default_factory=dict)
    entity_variants: dict[str, str] = field(default_factory=dict)
    entity_categories: dict[str, str] = field(default_factory=dict)
    detected_entities: dict[str, str] = field(default_factory=dict)
    mapping: dict[str, dict[str, str]] = field(default_factory=dict)
    counters: Counter[str] = field(default_factory=Counter)

    def _token(self, category: str, original: str) -> str:
        # One short code format (letter-digest) everywhere: filenames, folders,
        # document bodies, and PDFs. A single canonical form per entity means
        # there is no separate "compact" rendering to keep in sync elsewhere.
        key = f"{category}:{original}"
        if key in self.mapping:
            return self.mapping[key]["replacement"]
        labels = {
            "email": "person", "phone": "phone", "sin": "id", "card": "card",
            "ip": "ip", "organization": "company", "person": "person", "label": "label",
        }
        label = labels.get(category, category)
        letter = CODE_LETTER.get(label, label[0].upper())
        digest = hmac.new(self.secret, key.encode(), hashlib.sha256).hexdigest()
        used = {entry["replacement"] for entry in self.mapping.values()}
        length = 4
        replacement = f"{letter}-{digest[:length]}"
        while replacement in used and length < len(digest):
            length += 1
            replacement = f"{letter}-{digest[:length]}"
        self.mapping[key] = {"category": category, "original": original, "replacement": replacement}
        self.counters[category] += 1
        return replacement

    def _record_entity(self, original: str, replacement: str) -> None:
        key = f"entity:{original}"
        if key not in self.mapping:
            self.mapping[key] = {"category": "entity", "original": original, "replacement": replacement}
            self.counters["entity"] += 1

    def register_detected_entity(self, original: str, category: str) -> None:
        """Add an NER discovery while keeping variations of one name consistent."""
        key = normalise(original)
        if len(key) < 3 or key in self.configured_entities:
            return
        if key in self.detected_entities:
            self.entity_variants.setdefault(original, self.detected_entities[key])
            return
        # A detection such as "Apple Inc" may just be a longer/shorter spelling
        # of an entity we already know about (configured, e.g. bare "Apple", or
        # detected earlier in this run). Reuse that label so the same
        # real-world entity is never split across two different pseudonyms --
        # including product names built on it, like "Google Cloud Platform"
        # containing configured "Google", even when NER mis-tags the longer
        # phrase's category (e.g. tagging "Google Cloud" as a person).
        # Matches against another *auto-detected* (not explicitly configured)
        # entity stay category-scoped, so an organization can't absorb an
        # unrelated auto-detected person's label (or vice versa) purely from a
        # coincidentally shared word; configured entities are human-vetted, so
        # a word-span match against one is trusted regardless of category.
        words = _word_tokens(original)
        for known_original, replacement in self.entity_variants.items():
            known_norm = normalise(known_original)
            known_category = self.entity_categories.get(known_norm, "any")
            if known_norm not in self.configured_entities and known_category not in (category, "any"):
                continue
            if _shares_full_word_span(words, _word_tokens(known_original)):
                self.entity_variants[original] = replacement
                self.entity_categories.setdefault(key, known_category)
                return
        replacement = self._token(category, key)
        self.detected_entities[key] = replacement
        self.entity_categories[key] = category
        self.entity_variants[original] = replacement

    def find_entity_matches(self, text: str) -> dict[str, str]:
        """Return each literal substring of `text` matching a known entity,
        mapped to its replacement. Uses the same flexible word-boundary
        matching as replace_text (so punctuation variants such as a
        typographic apostrophe still match) but returns the exact text as
        written, so an exact-match consumer like PDF page search can find it.
        """
        matches: dict[str, str] = {}
        if not text:
            return matches
        for original, target in sorted(self.entity_variants.items(), key=lambda x: len(x[0]), reverse=True):
            words = _word_tokens(original)
            if not words:
                continue
            pattern = r"(?<![A-Za-z0-9])" + r"[\W_]*".join(map(re.escape, words)) + r"(?![A-Za-z0-9])"
            for m in re.finditer(pattern, text, flags=re.IGNORECASE):
                matches.setdefault(m.group(0), target)
        return matches

    def replace_text(self, value: str) -> str:
        if not isinstance(value, str) or not value:
            return value
        # Longest first prevents a short configured variation consuming a longer one.
        for original, target in sorted(self.entity_variants.items(), key=lambda x: len(x[0]), reverse=True):
            # Match punctuation/spacing variants by comparing word-like fragments.
            words = re.findall(r"[a-z0-9]+", original.casefold())
            if not words:
                continue
            pattern = r"(?<![A-Za-z0-9])" + r"[\W_]*".join(map(re.escape, words)) + r"(?![A-Za-z0-9])"
            def entity_sub(match: re.Match[str]) -> str:
                self._record_entity(match.group(0), target)
                return target
            value = re.sub(pattern, entity_sub, value, flags=re.IGNORECASE)
        value = EMAIL.sub(lambda m: self._token("email", m.group(0)), value)
        value = SIN.sub(lambda m: self._token("sin", m.group(0)), value)
        value = CARD.sub(lambda m: self._token("card", m.group(0)), value)
        value = IP.sub(lambda m: self._token("ip", m.group(0)), value)
        return PHONE.sub(lambda m: self._token("phone", m.group(0)), value)


def load_config(path: Path | None, p: Pseudonymizer) -> None:
    """Load configured company and personal-name variations."""
    payload: dict[str, Any] = {}
    if path:
        payload = json.loads(path.read_text(encoding="utf-8"))
    # Sections are separate only to make the config easy to audit. All use the
    # same replacement engine, including when text occurs in a path or title.
    # The JSON key (e.g. "Company Limited Three") is never the visible output;
    # it is only the seed for a short, stable code shared by every variant
    # listed under it, so the same entity looks identical everywhere.
    section_meta = {
        "companies": ("organization", "organization"),
        "people": ("person", "person"),
        "first_names": ("person", "person"),
        "last_names": ("person", "person"),
        "entities": ("label", "any"),
    }
    for section in ("entities", "companies", "people", "first_names", "last_names"):
        entries = payload.get(section, {})
        if not isinstance(entries, dict):
            raise ValueError(f"config '{section}' must be an object: replacement -> [known variations]")
        token_category, merge_category = section_meta[section]
        for label, variants in entries.items():
            if not isinstance(variants, list):
                raise ValueError(f"config '{section}.{label}' must be an array of variations")
            replacement = p._token(token_category, str(label))
            for variant in [label, *variants]:
                norm = normalise(str(variant))
                p.configured_entities[norm] = replacement
                p.entity_variants[str(variant)] = replacement
                p.entity_categories[norm] = merge_category
    # Do not guess that an arbitrary directory name is PII. Only configured
    # values are renamed, preserving the source structure exactly otherwise.


def extract_text_for_discovery(path: Path) -> str:
    """Read text safely for local entity discovery; failures simply yield no text."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            import pymupdf
            with pymupdf.open(path) as pdf:
                chunks = []
                for page in pdf:
                    chunks.append(page.get_text("text"))
                    if page.get_images(full=True):
                        pixmap = page.get_pixmap(matrix=pymupdf.Matrix(2, 2), alpha=False)
                        from PIL import Image
                        chunks.append(ocr_text(Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)))
                return "\n".join(chunks)
        if ext == ".docx":
            from docx import Document
            doc = Document(path)
            chunks = [paragraph.text for paragraph in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    chunks.extend(cell.text for cell in row.cells)
            return "\n".join(chunks)
        if ext in {".xlsx", ".xlsm"}:
            import openpyxl
            book = openpyxl.load_workbook(path, read_only=True, data_only=False, keep_vba=ext == ".xlsm")
            chunks: list[str] = []
            for sheet in book.worksheets:
                chunks.append(sheet.title)
                for row in sheet.iter_rows(values_only=True):
                    chunks.extend(str(cell) for cell in row if isinstance(cell, str))
            book.close()
            return "\n".join(chunks)
        if ext == ".pptx":
            from pptx import Presentation
            deck = Presentation(path)
            chunks = []
            for slide in deck.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        chunks.append(shape.text)
            return "\n".join(chunks)
        if ext in IMAGE_TYPES:
            from PIL import Image
            return ocr_text(Image.open(path))
        if ext in {".csv", ".tsv", ".txt", ".md", ".json", ".xml", ".html", ".htm"}:
            return path.read_text(encoding="utf-8", errors="replace")
    except Exception:
        return ""
    return ""


def ocr_text(image: Any) -> str:
    pytesseract = get_ocr_engine()
    return pytesseract.image_to_string(image)


def get_ocr_engine() -> Any:
    try:
        import pytesseract
        if not shutil.which("tesseract") and Path("/opt/homebrew/bin/tesseract").is_file():
            pytesseract.pytesseract.tesseract_cmd = "/opt/homebrew/bin/tesseract"
        pytesseract.get_tesseract_version()
        return pytesseract
    except (ImportError, OSError, RuntimeError) as exc:
        raise RuntimeError(
            "image OCR requires pytesseract and Tesseract. Install Python requirements and run 'brew install tesseract'"
        ) from exc


def ocr_redaction_boxes(image: Any, p: Pseudonymizer) -> list[tuple[int, int, int, int]]:
    """Return OCR word/phrase rectangles for known entity and pattern matches."""
    try:
        pytesseract = get_ocr_engine()
        from pytesseract import Output
        data = pytesseract.image_to_data(image, output_type=Output.DICT)
    except (ImportError, OSError, RuntimeError) as exc:
        raise RuntimeError(
            "image OCR requires pytesseract and Tesseract. Install Python requirements and run 'brew install tesseract'"
        ) from exc
    words: list[tuple[str, int, int, int, int]] = []
    for index, raw in enumerate(data["text"]):
        text = raw.strip()
        if text:
            words.append((text, int(data["left"][index]), int(data["top"][index]), int(data["width"][index]), int(data["height"][index])))
    matches: list[tuple[int, int]] = []
    # Longest configured/discovered entity phrases first, so an occurrence of
    # "Apple Inc." is redacted once rather than as both "Apple" and "Apple Inc.".
    candidates = sorted(p.entity_variants, key=len, reverse=True)
    normal_words = [normalise(word[0]) for word in words]
    for candidate in candidates:
        parts = [normalise(part) for part in re.findall(r"[A-Za-z0-9]+", candidate)]
        if not parts:
            continue
        for start in range(len(words) - len(parts) + 1):
            if normal_words[start:start + len(parts)] == parts:
                end = start + len(parts)
                if not any(start >= left and end <= right for left, right in matches):
                    matches.append((start, end))
    # Detect single-token structured PII such as email addresses as OCR permits.
    for index, (word, *_box) in enumerate(words):
        if p.replace_text(word) != word and not any(left <= index < right for left, right in matches):
            matches.append((index, index + 1))
    boxes: list[tuple[int, int, int, int]] = []
    for start, end in matches:
        selected = words[start:end]
        x0 = min(item[1] for item in selected); y0 = min(item[2] for item in selected)
        x1 = max(item[1] + item[3] for item in selected); y1 = max(item[2] + item[4] for item in selected)
        boxes.append((x0, y0, x1, y1))
    return boxes


def clean_image(src: Path, dest: Path, p: Pseudonymizer) -> None:
    from PIL import Image, ImageDraw
    with Image.open(src) as original:
        image = original.convert("RGBA") if original.mode == "P" else original.copy()
    for x0, y0, x1, y1 in ocr_redaction_boxes(image, p):
        ImageDraw.Draw(image).rectangle((x0, y0, x1, y1), fill="white")
    dest.parent.mkdir(parents=True, exist_ok=True)
    # Saving a newly-created image removes EXIF and other source metadata.
    if dest.suffix.lower() in {".jpg", ".jpeg"} and image.mode == "RGBA":
        image = image.convert("RGB")
    image.save(dest)


def discover_entities(source: Path, p: Pseudonymizer) -> None:
    """Discover PERSON and ORG names locally before transforming any files."""
    try:
        import spacy
        nlp = spacy.load("en_core_web_sm", disable=["tagger", "parser", "lemmatizer", "textcat"])
    except (ImportError, OSError) as exc:
        raise RuntimeError(
            "automatic name detection requires spaCy and its English model; run "
            "'.venv/bin/python -m spacy download en_core_web_sm'"
        ) from exc
    # The input convention is one company per top-level folder. Unlike a bare
    # word in document text, this is an explicit, high-confidence organization
    # hint and must be anonymized even without a legal suffix such as "Ltd".
    for child in source.iterdir():
        if child.is_dir() and not child.name.startswith("."):
            p.register_detected_entity(child.name, "organization")
    samples: list[str] = []
    for item in source.rglob("*"):
        if item.is_file():
            samples.append(" ".join(item.relative_to(source).with_suffix("").parts))
            text = extract_text_for_discovery(item)
            if text:
                samples.append(text)
    candidates: dict[tuple[str, str], dict[str, Any]] = {}
    for document in nlp.pipe(samples, batch_size=16):
        for entity in document.ents:
            if entity.label_ not in {"ORG", "PERSON"}:
                continue
            original = entity.text.strip()
            key = (entity.label_, normalise(original))
            if len(key[1]) < 3:
                continue
            candidate = candidates.setdefault(key, {"forms": [], "count": 0})
            candidate["count"] += 1
            if original not in candidate["forms"]:
                candidate["forms"].append(original)

    for (label, _), candidate in candidates.items():
        original = candidate["forms"][0]
        words = re.findall(r"[A-Za-z]+", original)
        if label == "PERSON":
            # Single-token labels (for example, a capitalized heading) are too
            # ambiguous to redact safely without an explicit mapping.
            eligible = len(words) >= 2 and all(word[0].isupper() for word in words)
            category = "person"
        else:
            # Unknown organizations require a legal/company cue. Bare names such
            # as "Apple" should be supplied in the explicit config instead.
            eligible = bool(COMPANY_CUE.search(original))
            category = "organization"
        if eligible:
            for form in candidate["forms"]:
                p.register_detected_entity(form, category)


def sanitized_name(name: str, p: Pseudonymizer) -> str:
    stem, suffix = os.path.splitext(name)
    return p.replace_text(stem) + suffix


def clean_docx(src: Path, dest: Path, p: Pseudonymizer) -> None:
    from docx import Document
    doc = Document(src)
    def paragraphs(items: Iterable[Any]) -> None:
        for para in items:
            for run in para.runs:
                run.text = p.replace_text(run.text)
    paragraphs(doc.paragraphs)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                paragraphs(cell.paragraphs)
    for section in doc.sections:
        paragraphs(section.header.paragraphs); paragraphs(section.footer.paragraphs)
    props = doc.core_properties
    for attr in ("author", "title", "subject", "comments", "keywords", "last_modified_by", "category", "identifier"):
        setattr(props, attr, p.replace_text(getattr(props, attr, "") or ""))
    props.created = props.modified = datetime.now(timezone.utc)
    dest.parent.mkdir(parents=True, exist_ok=True); doc.save(dest)


def clean_xlsx(src: Path, dest: Path, p: Pseudonymizer) -> None:
    import openpyxl
    wb = openpyxl.load_workbook(src, keep_vba=src.suffix.lower() == ".xlsm")
    for sheet in wb.worksheets:
        sheet.title = p.replace_text(sheet.title)[:31]
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and not cell.value.startswith("="):
                    cell.value = p.replace_text(cell.value)
                if cell.comment:
                    cell.comment.text = p.replace_text(cell.comment.text)
    props = wb.properties
    for attr in ("creator", "title", "subject", "description", "keywords", "lastModifiedBy"):
        setattr(props, attr, p.replace_text(getattr(props, attr, "") or ""))
    dest.parent.mkdir(parents=True, exist_ok=True); wb.save(dest)


def clean_pptx(src: Path, dest: Path, p: Pseudonymizer) -> None:
    from pptx import Presentation
    deck = Presentation(src)
    for slide in deck.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs: run.text = p.replace_text(run.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells: cell.text = p.replace_text(cell.text)
    props = deck.core_properties
    for attr in ("author", "title", "subject", "comments", "keywords", "last_modified_by", "category"):
        setattr(props, attr, p.replace_text(getattr(props, attr, "") or ""))
    dest.parent.mkdir(parents=True, exist_ok=True); deck.save(dest)


def clean_pdf(src: Path, dest: Path, p: Pseudonymizer) -> int:
    import pymupdf as fitz
    from PIL import Image
    pdf = fitz.open(src)
    replacements: dict[str, str] = {}
    image_only_pages = 0
    for page in pdf:
        text = page.get_text("text")
        if not text.strip() and page.get_images(full=True):
            image_only_pages += 1
        # Map detected literal strings to their pseudonyms. PDF redaction is exact,
        # so entity matches must be found via the same flexible word-boundary regex
        # replace_text uses (not a raw substring check), and recorded using the
        # exact text found on the page -- a document may render a configured
        # variant with different punctuation (e.g. a typographic apostrophe).
        for pattern in (EMAIL, SIN, CARD, IP, PHONE):
            for match in pattern.finditer(text):
                original = match.group(0)
                replacement = p.replace_text(original)
                if original != replacement: replacements[original] = replacement
        for original, replacement in p.find_entity_matches(text).items():
            replacements[original] = replacement
            p._record_entity(original, replacement)
    for page in pdf:
        occupied: list[Any] = []
        # A shorter variant ("Apple") can occur inside a longer variant
        # ("Apple Inc."). Redact the longest one first and never draw two labels
        # over the same PDF text area.
        for original, replacement in sorted(replacements.items(), key=lambda item: len(item[0]), reverse=True):
            for rect in page.search_for(original):
                if any(rect.intersects(previous) for previous in occupied):
                    continue
                page.add_redact_annot(
                    rect,
                    text=replacement,
                    fill=(1, 1, 1),
                    text_color=(0, 0, 0),
                    fontsize=5,
                )
                occupied.append(rect)
        # OCR all rendered pages that contain images. The redaction operation is
        # applied to image pixels, not merely drawn over the top, so recovered
        # image streams cannot reveal the matched text.
        if page.get_images(full=True):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
            scale_x = page.rect.width / pixmap.width
            scale_y = page.rect.height / pixmap.height
            for x0, y0, x1, y1 in ocr_redaction_boxes(image, p):
                rect = fitz.Rect(x0 * scale_x, y0 * scale_y, x1 * scale_x, y1 * scale_y)
                if any(rect.intersects(previous) for previous in occupied):
                    continue
                page.add_redact_annot(rect, fill=(1, 1, 1))
                occupied.append(rect)
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)
    meta = pdf.metadata or {}
    for key, value in list(meta.items()): meta[key] = p.replace_text(value or "")
    meta["author"] = "Anonymized"; meta["producer"] = "PII Anonymizer"
    pdf.set_metadata(meta)
    dest.parent.mkdir(parents=True, exist_ok=True); pdf.save(dest, garbage=4, deflate=True)
    pdf.close()
    return image_only_pages


def clean_text(src: Path, dest: Path, p: Pseudonymizer) -> None:
    if src.suffix.lower() == ".json":
        def clean_json(value: Any) -> Any:
            if isinstance(value, dict):
                return {p.replace_text(str(key)): clean_json(item) for key, item in value.items()}
            if isinstance(value, list):
                return [clean_json(item) for item in value]
            return p.replace_text(value) if isinstance(value, str) else value
        try:
            content = json.loads(src.read_text(encoding="utf-8-sig"))
        except json.JSONDecodeError:
            # Preserve malformed JSON as text, while still sanitizing it.
            pass
        else:
            dest.write_text(json.dumps(clean_json(content), indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
            return
    if src.suffix.lower() in {".csv", ".tsv"}:
        dialect = csv.excel_tab if src.suffix.lower() == ".tsv" else csv.excel
        with src.open("r", encoding="utf-8-sig", newline="") as fin, dest.open("w", encoding="utf-8", newline="") as fout:
            writer = csv.writer(fout, dialect=dialect)
            for row in csv.reader(fin, dialect=dialect): writer.writerow([p.replace_text(x) for x in row])
        return
    content = src.read_text(encoding="utf-8", errors="replace")
    dest.write_text(p.replace_text(content), encoding="utf-8")


def process_file(src: Path, dest: Path, p: Pseudonymizer, allow_copy: bool) -> str:
    ext = src.suffix.lower()
    if ext == ".docx": clean_docx(src, dest, p)
    elif ext in {".xlsx", ".xlsm"}: clean_xlsx(src, dest, p)
    elif ext == ".pptx": clean_pptx(src, dest, p)
    elif ext == ".pdf":
        image_only_pages = clean_pdf(src, dest, p)
        if image_only_pages:
            return f"anonymized; OCR applied to {image_only_pages} image-only page(s); review required"
    elif ext in IMAGE_TYPES:
        clean_image(src, dest, p)
    elif ext in {".csv", ".tsv", ".txt", ".md", ".json", ".xml", ".html", ".htm"}:
        dest.parent.mkdir(parents=True, exist_ok=True); clean_text(src, dest, p)
    elif allow_copy:
        dest.parent.mkdir(parents=True, exist_ok=True); shutil.copy2(src, dest); return "copied-unprocessed"
    else: return "skipped-unsupported"
    return "anonymized"


def write_mapping(path: Path, p: Pseudonymizer) -> None:
    payload = {"generated_at": datetime.now(timezone.utc).isoformat(), "entries": list(p.mapping.values())}
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("source", type=Path, help="root directory containing original attachments")
    ap.add_argument("output", type=Path, help="new output directory (must not be inside source)")
    ap.add_argument("--config", type=Path, help="JSON file of replacement entities and known variations")
    ap.add_argument("--copy-unsupported", action="store_true", help="copy unsupported files unchanged (not anonymized)")
    ap.add_argument("--secret", help="stable secret for repeatable pseudonyms; prefer PII_ANON_SECRET env var")
    ap.add_argument("--overwrite-output", action="store_true", help="delete and recreate the specified output directory before processing")
    ap.add_argument("--no-auto-detect-names", action="store_true", help="disable local detection of unknown organization and person names")
    args = ap.parse_args()
    source, output = args.source.resolve(), args.output.resolve()
    if not source.is_dir(): ap.error("source must be an existing directory")
    if output == source or source in output.parents: ap.error("output must not be source or a child of source")
    if output.exists() and any(output.iterdir()):
        if not args.overwrite_output:
            ap.error("output must be empty or not exist (use --overwrite-output to replace it)")
        shutil.rmtree(output)
    secret = (args.secret or os.getenv("PII_ANON_SECRET") or os.urandom(32).hex()).encode()
    p = Pseudonymizer(secret)
    load_config(args.config, p)
    if not args.no_auto_detect_names:
        try:
            discover_entities(source, p)
        except RuntimeError as exc:
            ap.error(str(exc))
    audit: list[dict[str, str]] = []
    claimed_files: dict[Path, Path] = {}

    # Recreate even empty directories. Each relative segment is anonymized, so
    # hierarchy is unchanged while PII in folder titles is replaced.
    for directory in sorted((x for x in source.rglob("*") if x.is_dir()), key=lambda x: len(x.parts)):
        rel = directory.relative_to(source)
        output.joinpath(*(sanitized_name(part, p) for part in rel.parts)).mkdir(parents=True, exist_ok=True)

    for src in source.rglob("*"):
        if not src.is_file(): continue
        rel = src.relative_to(source)
        out_parts = [sanitized_name(part, p) for part in rel.parts]
        dest = output.joinpath(*out_parts)
        previous = claimed_files.get(dest)
        if previous and previous != src:
            status = f"failed: anonymized path collision with {previous.relative_to(source)}"
        else:
            claimed_files[dest] = src
            try: status = process_file(src, dest, p, args.copy_unsupported)
            except Exception as exc: status = f"failed: {type(exc).__name__}: {exc}"
        audit.append({"source": str(rel), "output": str(dest.relative_to(output)), "status": status})
    output.mkdir(parents=True, exist_ok=True)
    (output / "anonymization_audit.json").write_text(json.dumps({"files": audit, "counts": p.counters}, indent=2), encoding="utf-8")
    write_mapping(output / "reverse_mapping.json", p)
    failures = [x for x in audit if x["status"].startswith("failed") or x["status"].startswith("skipped")]
    print(f"Processed {len(audit)} files; {len(failures)} require review. Output: {output}")
    return 1 if failures else 0

if __name__ == "__main__":
    raise SystemExit(main())
